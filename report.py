from pptx import Presentation
from pptx.util import Inches
import os
import time
import warnings
warnings.filterwarnings("ignore", category=FutureWarning)

def create_report(image_path, result, filename):
    # Create a PowerPoint presentation object
    prs = Presentation()

    # Add title slide
    slide_layout = prs.slide_layouts[0]  # 0 is the title slide layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Plastic Detection Report"
    subtitle.text = f"Analysis for {filename}"

    # Add slide for image
    slide_layout = prs.slide_layouts[5]  # 5 is a blank slide layout
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.add_picture(image_path, Inches(0.5), Inches(1), height=Inches(4.5))

    # Add slide for analysis result
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # 1 is the Title and Content layout
    title = slide.shapes.title
    title.text = "Analysis Result"

    content = slide.shapes.placeholders[1]
    content.text = result

    # Ensure the reports folder exists
    if not os.path.exists('reports'):
        os.makedirs('reports')

    # Create a unique filename for the report
    timestamp = int(time.time())
    report_filename = f"{timestamp}_{filename.split('.')[0]}_report.pptx"

    # Save the presentation to a file in the 'reports' folder
    report_path = os.path.join('reports', report_filename)
    prs.save(report_path)

    return report_path


def create_report_2(image_path, result, filename):
    # Create a PowerPoint presentation object
    prs = Presentation()

    # Add title slide
    slide_layout = prs.slide_layouts[0]  # 0 is the title slide layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Plastic Detection Report"
    subtitle.text = f"Analysis for {filename}"

    # Add slide for image
    slide_layout = prs.slide_layouts[5]  # 5 is a blank slide layout
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.add_picture(image_path, Inches(0.5), Inches(1), height=Inches(4.5))

    # Add slide for analysis description
    slide_layout = prs.slide_layouts[1]  # Title and Content layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Analysis Overview"
    content = slide.shapes.placeholders[1]
    content.text = result["description"]

    # Add slides for each object detected
    for obj in result['objects']:
        slide_layout = prs.slide_layouts[1]  # Title and Content layout
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = f"Object: {obj['type']}"

        content = slide.shapes.placeholders[1]
        content.text = (
            f"Shape (Aspect Ratio): {obj['shape']}\n"
            f"Color: {obj['color']}\n"
            f"Confidence: {obj['confidence']:.2f}"
        )

    # Ensure the reports folder exists
    if not os.path.exists('reports'):
        os.makedirs('reports')

    # Create a unique filename for the report
    timestamp = int(time.time())
    report_filename = f"{timestamp}_{filename.split('.')[0]}_report.pptx"

    # Save the presentation to a file in the 'reports' folder
    report_path = os.path.join('reports', report_filename)
    prs.save(report_path)

    return report_path

from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
import os
import time

def create_report_3(image_path, result, filename):
    # Ensure the reports folder exists
    if not os.path.exists('reports'):
        os.makedirs('reports')

    # Create a unique filename for the report
    timestamp = int(time.time())
    report_filename = f"{timestamp}_{filename.split('.')[0]}_report.pdf"
    report_path = os.path.join('reports', report_filename)

    # Create a canvas for the PDF
    c = canvas.Canvas(report_path, pagesize=letter)
    width, height = letter

    # Add title page
    c.setFont("Helvetica-Bold", 24)
    c.drawString(100, height - 100, "Plastic Detection Report")
    c.setFont("Helvetica", 12)
    c.drawString(100, height - 130, f"Analysis for {filename}")
    c.showPage()  # End the current page

    # Add image to PDF
    c.drawImage(image_path, 100, height - 450, width=400, height=300)
    c.showPage()  # End the current page

    # Add analysis description slide
    c.setFont("Helvetica-Bold", 18)
    c.drawString(100, height - 100, "Analysis Overview")
    c.setFont("Helvetica", 12)
    c.drawString(100, height - 130, result["description"])
    c.showPage()  # End the current page

    # Add slides for each object detected
    for obj in result['objects']:
        c.setFont("Helvetica-Bold", 18)
        c.drawString(100, height - 100, f"Object: {obj['type']}")
        c.setFont("Helvetica", 12)
        c.drawString(100, height - 130, f"Shape (Aspect Ratio): {obj['shape']}")
        c.drawString(100, height - 150, f"Color: {obj['color']}")
        c.drawString(100, height - 170, f"Confidence: {obj['confidence']:.2f}")
        c.showPage()  # End the current page

    # Save the PDF file
    c.save()

    return report_path

def check_for_pyrolysis(object_type):
    """
    Check if the detected object is useful for the pyrolysis process based on its type.
    
    Args:
    - object_type (str): The type of the detected object (e.g., "plastic", "paper", "metal").
    
    Returns:
    - bool: True if the object is useful for pyrolysis, False otherwise.
    """
    # List of object types useful for pyrolysis (plastic types commonly used for pyrolysis)
    pyrolysis_compatible_types = [
        "polyethylene", "polypropylene", "polystyrene", "polyvinyl chloride", "polyethylene terephthalate",
        "bottle", "cup", "cover", "wrapper", "bag", "can", "container", "cell phone", "remote", "toilet"
    ]

    # If the object type matches one of the compatible types, return True (useful for pyrolysis)
    if object_type.lower() in pyrolysis_compatible_types:
        return True
    else:
        return False


from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
import os
import time
from reportlab.lib import colors

def create_report_4(image_path, result, filename):
    # Ensure the reports folder exists
    if not os.path.exists('reports'):
        os.makedirs('reports')

    # Create a unique filename for the report
    timestamp = int(time.time())
    report_filename = f"{timestamp}_{filename.split('.')[0]}_report.pdf"
    report_path = os.path.join('reports', report_filename)

    # Create a canvas for the PDF
    c = canvas.Canvas(report_path, pagesize=letter)
    width, height = letter

    # **First page**: Display logo (Favicon.png)
    c.drawImage(r"C:\Users\hp\Desktop\Plastic_detection\static\images\Favicon.png", 100, height - 150, width=200, height=150)
    c.setFont("Helvetica-Bold", 24)
    c.drawString(100, height - 250, "Plastic Detection Report")
    c.setFont("Helvetica", 12)
    c.drawString(100, height - 280, f"Analysis for {filename}")
    c.showPage()  # End first page

    # **Second page**: Display the uploaded image
    c.setFont("Helvetica-Bold", 18)
    c.drawString(100, height - 100, "This is the image we have analyzed:")
    c.drawImage(image_path, 100, height - 450, width=400, height=300)
    c.showPage()  # End second page

    # **Third page**: Number of objects and their names
    c.setFont("Helvetica-Bold", 18)
    c.drawString(100, height - 100, "Number of Objects Detected and Their Names:")
    y_position = height - 130  # Starting position for list
    for obj in result['objects']:
        c.setFont("Helvetica", 12)
        c.drawString(100, y_position, f"- {obj['type']}")
        y_position -= 20  # Adjust the y-position for each object
    c.showPage()  # End third page

    # **From the 4th page onwards**: Description of each object
    for obj in result['objects']:
        c.setFont("Helvetica-Bold", 18)
        c.drawString(100, height - 100, f"Object: {obj['type']}")

        # # Display the object image (small size, 1/5th of the page size)
        # small_image_path = obj['image']  # You can provide the path to a cropped or resized image
        # c.drawImage(small_image_path, width - 200, height - 150, width=100, height=100)

        # Display object attributes
        c.setFont("Helvetica", 12)
        c.drawString(100, height - 180, f"Shape: {obj['shape']}")
        c.drawString(100, height - 200, f"Color: {obj['color']}")
        c.drawString(100, height - 220, f"Confidence: {obj['confidence']:.2f}")
        #c.drawString(100, height - 240, f"Useful for Pyrolysis: {'Yes' if obj['useful_for_pyrolysis'] else 'No'}")
        useful_for_pyrolysis = check_for_pyrolysis(obj['type'])
        c.drawString(100, height - 240, f"Useful for Pyrolysis: {'Yes' if useful_for_pyrolysis else 'No'}")


        # Check if there are more attributes to display
        if 'additional_attributes' in obj:
            y_position = height - 260
            for attribute, value in obj['additional_attributes'].items():
                c.drawString(100, y_position, f"{attribute}: {value}")
                y_position -= 20

        c.showPage()  # End page for the object

    # Save the PDF file
    c.save()

    return report_path
